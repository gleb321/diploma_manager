import os
import smtplib
import qrcode
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.application import MIMEApplication
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Inches


def send_file(filename, login, password, email):
    try:
        msg = MIMEMultipart()
        msg['From'] = login
        msg['To'] = email
        msg['Subject'] = "Skillbox diploma"
        with open(filename, "rb") as data:
            attached_file = MIMEApplication(data.read())
            attached_file.add_header('Content-Disposition', 'attachment', filename = filename)
            msg.attach(attached_file)
        server = smtplib.SMTP('smtp.gmail.com', 587)
        server.starttls()
        server.login(msg['From'], password)
        server.sendmail(msg['From'], msg['To'], msg.as_string())
        server.quit()
    except Exception as ex:
        raise ex


def create_diploma(diploma_name = "diploma.pptx", template = "template.pptx", name = "", surname = "", course = "", date = "", link = ""):
    try:
        presentation = Presentation(template)
        slide = presentation.slides[0]
        for shape in slide.shapes:
            if shape.has_text_frame:
                if (shape.text_frame.paragraphs[0].text == "Initials"):
                    shape.text_frame.paragraphs[0].font.size = Pt(18);
                    shape.text_frame.paragraphs[0].text = f"{name} {surname}"
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                if (shape.text_frame.paragraphs[0].text == "Course"):
                    shape.text_frame.paragraphs[0].font.size = Pt(18);
                    shape.text_frame.paragraphs[0].text = course
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                if (shape.text_frame.paragraphs[0].text == "Date"):
                    shape.text_frame.paragraphs[0].font.size = Pt(14);
                    shape.text_frame.paragraphs[0].text = f'{date}'
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

                if (shape.text_frame.paragraphs[0].text == "Number"):
                    try: 
                        with open("number.txt", 'r') as reader:
                            number = int(reader.read())
                            number = number + 1
                    except Exception as ex:
                        raise Exception("Не удалось получить номер диплома")

                    if (6 - len(str(number)) >= 0):
                        zeros = (6 - len(str(number))) * "0"
                    else:
                        zeros = ""
                    shape.text_frame.paragraphs[0].font.size = Pt(14);
                    shape.text_frame.paragraphs[0].text = f'№ {zeros + str(number)}'
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

                    try: 
                        with open("number.txt", 'w') as writer:
                            writer.write(str(number))
                    except Exception as ex:
                        raise Exception("Не удалось записать номер диплома")

        qrcode.make(link).save("qr.png")
        slide.shapes.add_picture("qr.png", left = Inches(3.4), top = Inches(9.06), width = Inches(1), height = Inches(1))
        presentation.save(diploma_name)
        os.system(f"soffice --headless --convert-to pdf {diploma_name}")
    except Exception as ex:
        raise ex
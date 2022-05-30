import json
import _thread
import smtplib
import hashlib
import jwt
import pytz
import qrcode
import os
import time
import sqlite3 as sq
import pywaves as pw
from datetime import datetime
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.application import MIMEApplication
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Inches
from threading import Lock
from flask import Flask, request, jsonify
from flask_cors import CORS, cross_origin
from config import gmail_login, gmail_password, waves_private_key, host, port, jwt_key

db = Lock()

app = Flask(__name__)
cors = CORS(app)
app.config['CORS_HEADERS'] = 'Content-Type'


def parse_request(json_data):
    try:
        return json.loads(json_data)
    except Exception as ex:
        print("Не удалось распарсить json:")
        print(ex)
        raise ex


def create_db():
    try:
        with sq.connect("data.db") as con:
            cur = con.cursor()

            cur.execute("""CREATE TABLE IF NOT EXISTS diplomas (
                name text not null,
                surname text not null,
                email text not null,
                course text not null,
                course_id integer not null,
                direction text not null,
                portfolio text not null,
                hashcode text not null
                )""")

            print("Таблица requests успешно создана")

            cur.execute("""CREATE TABLE IF NOT EXISTS transactions (
                hashcode text not null,
                transaction_id text not null,
                date text not null
                )""")

            print("Таблица transactions успешно создана")

            cur.execute("""CREATE TABLE IF NOT EXISTS users (
                login text not null,
                password text not null,
                role text not null
                )""")

            print("Таблица users успешно создана")

            cur.close()
    except Exception as ex:
        print("Не удалось создать таблицу:")
        print(ex)


def add_to_db(data, case):
    db.acquire()
    try:
        with sq.connect("data.db") as con:
            cur = con.cursor()
            if case == "diploma":
                cur.execute("""INSERT INTO diplomas (name, surname, email, course, course_id, direction, portfolio, hashcode) 
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?)""", data)
                print("Данные о запросе успешно сохранены")
            elif case == "transaction":
                cur.execute("""INSERT INTO transactions (hashcode, transaction_id, date) VALUES (?, ?, ?)""", data)
                print("Данные о транзакции успешно сохранены")
            elif case == "user":
                cur.execute("""INSERT INTO users (login, password, role) VALUES (?, ?, ?)""", data)
                print("Данные о пользователе успешно сохранены")
            cur.close()
    except Exception as ex:
        if case == "diploma":
            print("Не удалось сохранить данные о запросе:")
        elif case == "transaction":
            print("Не удалось сохранить данные о транзакции:")
        elif case == "user":
            print("Не удалось сохранить данные о пользователе:")
        if db.locked():
            db.release()
        print(ex)
        raise ex

    if db.locked():
        db.release()


def get_user_data(login):
    db.acquire()
    client = None
    try:
        with sq.connect("data.db") as con:
            cur = con.cursor()
            users = cur.execute("""SELECT * FROM users WHERE login = ?""", (login, ))
            for user in users:
                client = user
            cur.close()
    except Exception as ex:
        print("Не удалось получить данные о пользователе:")
        print(ex)
        if db.locked():
            db.release()
        raise ex   

    if db.locked():
        db.release()

    if (client == None):
        raise Exception("Пользователь с таким логином не найден")

    return client[1], client[2]


def create_diploma(diploma_name = "diploma.pptx", template = "template.pptx", name = "", surname = "", course = "", date = "", link = ""):
    try:
        presentation = Presentation(template)
        slide = presentation.slides[0]
        for shape in slide.shapes:
            if shape.has_text_frame:
                if (shape.text_frame.paragraphs[0].text == "Initials"):
                    shape.text_frame.paragraphs[0].font.size = Pt(18);
                    shape.text_frame.paragraphs[0].text = f"{name} {surname}"
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                if (shape.text_frame.paragraphs[0].text == "Course"):
                    shape.text_frame.paragraphs[0].font.size = Pt(18);
                    shape.text_frame.paragraphs[0].text = course
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                if (shape.text_frame.paragraphs[0].text == "Date"):
                    shape.text_frame.paragraphs[0].font.size = Pt(14);
                    shape.text_frame.paragraphs[0].text = f'{date}'
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

                if (shape.text_frame.paragraphs[0].text == "Number"):
                    try: 
                        with open("number.txt", 'r') as reader:
                            number = int(reader.read())
                            number = number + 1
                    except Exception as ex:
                        print("Не удалось получить номер диплома:")
                        print(ex)
                        raise Exception("Не удалось получить номер диплома")

                    if (6 - len(str(number)) >= 0):
                        zeros = (6 - len(str(number))) * "0"
                    else:
                        zeros = ""
                    shape.text_frame.paragraphs[0].font.size = Pt(14);
                    shape.text_frame.paragraphs[0].text = f'№ {zeros + str(number)}'
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

                    try: 
                        with open("number.txt", 'w') as writer:
                            writer.write(str(number))
                    except Exception as ex:
                        print("Не удалось записать номер диплома:")
                        print(ex)
                        raise Exception("Не удалось записать номер диплома")

        qrcode.make(link).save("qr.png")
        slide.shapes.add_picture("qr.png", left = Inches(3.4), top = Inches(9.06), width = Inches(1), height = Inches(1))
        presentation.save(diploma_name)
        os.system(f"soffice --headless --convert-to pdf {diploma_name}")
        print("Диплом успешно создан")
    except Exception as ex:
        print("Не удалось создать диплом:")
        print(ex)
        raise ex


def send_file(filename = "diploma.pdf", login = gmail_login, password = gmail_password, email = gmail_login):
    try:
        msg = MIMEMultipart()
        msg['From'] = login
        msg['To'] = email
        msg['Subject'] = "Skillbox diploma"
        with open(filename, "rb") as data:
            attached_file = MIMEApplication(data.read())
            attached_file.add_header('Content-Disposition', 'attachment', filename = filename)
            msg.attach(attached_file)
        server = smtplib.SMTP('smtp.gmail.com', 587)
        server.starttls()
        server.login(msg['From'], password)
        server.sendmail(msg['From'], msg['To'], msg.as_string())
        server.quit()
        print("Диплом успешно отправлен")
    except Exception as ex:
        print("Не удалось отправить диплом:")
        print(ex)
        raise ex


def create_transaction(name, surname, course_id, date, hashcode):
    try:
        myAddress = pw.Address(privateKey = waves_private_key)
        data = [{'type': 'string', 'key': 'hashcode', 'value': hashcode}]
        transaction = myAddress.dataTransaction(data)
        print("Транзакция успешно совершена")
    except Exception as ex:
        print("Не удалось совершить транзакцию:")
        print(ex)
        raise ex

    add_to_db((hashcode, transaction['id'], date), 'transaction')
    transaction_id = transaction['id']
    return f'https://wavesexplorer.com/tx/{transaction_id}'


def get_diplomas(case = "all", email = ""):
    db.acquire()
    try:
        data = []
        with sq.connect("data.db") as con:
                cur = con.cursor()
                if (case == "all"):
                    diplomas = cur.execute("""SELECT surname, name, direction, course, email, portfolio, date, transaction_id
                         FROM diplomas JOIN transactions ON diplomas.hashcode = transactions.hashcode""")
                elif (case == "user"):
                     diplomas = cur.execute("""SELECT surname, name, direction, course, email, portfolio, date, transaction_id, course_id
                         FROM diplomas JOIN transactions ON diplomas.hashcode = transactions.hashcode
                         WHERE email = ?""", (email, ))
                for diploma in diplomas:
                    dict_data = {}
                    dict_data["surname"] = diploma[0]
                    dict_data["name"] = diploma[1]
                    dict_data["direction"] = diploma[2]
                    dict_data["course"] = diploma[3]
                    dict_data["email"] = diploma[4]
                    dict_data["portfolio"] = diploma[5]
                    dict_data["date"] = diploma[6]
                    dict_data["transaction"] = f'https://wavesexplorer.com/tx/{diploma[7]}'
                    if (case == "user"):
                        dict_data["course_id"] = diploma[8]
                    data.append(dict_data)
                cur.close()
        if db.locked():
            db.release()
        return str(json.dumps(data))
    except Exception as ex:
        print("Не удалось получить данные о дипломах:")
        print(ex)
        if db.locked():
            db.release()
        raise ex


@app.route('/request', methods=['GET', 'POST'])
@cross_origin()
def get_request():
    try:
        login, token = request.args.get('login'), request.args.get('token')
        password, role = get_user_data(login)
        try:
            decoded_token = jwt.decode(token, f"{login}.{password}.{jwt_key}", algorithms=["HS256"])
            if (decoded_token["role"] != "admin"):
                raise Exception("Недостаточный уровень доступа")

            if (decoded_token["exp"] < int(time.time())):
                raise Exception("Время действия токена истекло")

        except Exception as ex:
            raise Exception(f"Некорректный токен")

        data = dict(request.get_json())
        date = datetime.now(pytz.timezone('Europe/Moscow')).date()
        print(data)
        print("Получен новый запрос")
        string_to_hash = f'{data["name"]}_{data["surname"]}_{data["course_id"]}'
        hashcode = hashlib.sha256(string_to_hash.encode()).hexdigest()
        link = create_transaction(data["name"], data["surname"], data["course_id"], date.strftime("%d/%m/%Y"), hashcode)
        create_diploma(name = data["name"], surname = data["surname"], course = data["course"], date = date.strftime("%d.%m.%Y"), link = link)
        add_to_db((*data.values(), hashcode), "diploma")
        send_file(email = data["email"])
        print("Запрос успешно обработан")
        return "Запрос успешно обработан"
    except Exception as ex:
        print("Не удалось обработать запрос:")
        print(ex)
        return f"Не удалось обработать запрос:\n{ex}", 400


@app.route('/auth', methods=['GET'])
@cross_origin()
def authorize():
    login, password = request.args.get('login'), request.args.get('password')
    try:
        user_password, user_role = get_user_data(login)
        if (password != user_password):
            raise Exception("Неверный пароль")

        minutes = 15
        jwt_token = jwt.encode({"sub": "auth", "user": login, "role": user_role, "exp": int(time.time()) + 60 * minutes},
             f"{login}.{password}.{jwt_key}", algorithm="HS256")
        print("Авторизация прошла успешно")
        return jwt_token
    except Exception as ex:
        print("Ошибка авторизации:")
        print(ex)
        return f"Ошибка авторизации:\n{ex}", 404


@app.route('/diplomas', methods=['GET'])
@cross_origin()
def get_all_diplomas():
    try:
        login, token = request.args.get('login'), request.args.get('token')
        password, role = get_user_data(login)
        try:
            decoded_token = jwt.decode(token, f"{login}.{password}.{jwt_key}", algorithms=["HS256"])
            if (decoded_token["exp"] < int(time.time())):
                raise Exception("Время действия токена истекло")

            if (decoded_token["role"] != "admin"):
                raise Exception("Недостаточный уровень доступа")

        except Exception as ex:
            raise Exception("Некорректный токен")

        return get_diplomas(case = "all")
    except Exception as ex:
        print("Не удалось получить информацию о дипломах:")
        print(ex)
        return f"Не удалось получить информацию о дипломах:\n{ex}", 400


@app.route('/user', methods=['GET'])
@cross_origin()
def get_user_diplomas():
    try:
        login, token = request.args.get('login'), request.args.get('token')
        password, role = get_user_data(login)
        try:
            decoded_token = jwt.decode(token, f"{login}.{password}.{jwt_key}", algorithms=["HS256"])
            if (decoded_token["exp"] < int(time.time())):
                raise Exception("Время действия токена истекло")

        except Exception as ex:
            raise Exception("Некорректный токен")

        return get_diplomas(case = "user", email = login)
    except Exception as ex:
        print("Не удалось получить информацию о дипломах пользователя:")
        print(ex)
        return f"Не удалось получить информацию о дипломах пользователя:\n{ex}", 400


if __name__ == '__main__':
    app.run(host = host, port = port)











